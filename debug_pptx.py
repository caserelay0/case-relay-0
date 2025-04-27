import os
import sys
import logging
import io
from pathlib import Path
import base64
from PIL import Image
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

def process_pptx_images(filepath):
    """Test function to process images from a PPTX file with robust error handling"""
    logger.info(f"Testing PPTX image extraction: {filepath}")
    
    images = []
    image_count = 0
    max_images = 10  # Process just a few images
    
    try:
        presentation = Presentation(filepath)
        
        # Process each slide
        for slide_num, slide in enumerate(presentation.slides[:5]):  # First 5 slides only
            logger.info(f"Processing slide {slide_num + 1}")
            
            # Process shapes recursively with robust error handling
            def process_shape(shape, shape_id="unknown"):
                nonlocal image_count
                
                if image_count >= max_images:
                    return
                
                try:
                    logger.debug(f"Processing shape type: {shape.shape_type} (id: {shape_id})")
                    
                    # Handle image shapes - most common type
                    if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # Validate image data is available
                            if not hasattr(shape, 'image') or not shape.image:
                                logger.warning(f"Shape has no image attribute or image is None")
                                return
                                
                            # Get image bytes
                            image_bytes = shape.image.blob
                            if not image_bytes:
                                logger.warning(f"Image blob is empty")
                                return
                                
                            # Successfully got image data
                            logger.info(f"Found image in shape {shape_id}, size: {len(image_bytes)} bytes")
                            
                            # Test image opening and conversion
                            try:
                                # Try opening the image
                                img_stream = io.BytesIO(image_bytes)
                                img = Image.open(img_stream)
                                logger.info(f"Image opened successfully: format={img.format}, mode={img.mode}, size={img.size}")
                                
                                # Try saving with different formats
                                for fmt in ['PNG', 'JPEG']:
                                    try:
                                        # Create a new buffer for each attempt
                                        test_buffer = io.BytesIO()
                                        
                                        # Convert to RGB if needed for JPEG
                                        if fmt == 'JPEG' and img.mode in ['RGBA', 'LA']:
                                            img_rgb = img.convert('RGB')
                                            img_rgb.save(test_buffer, format=fmt)
                                            logger.info(f"Successfully saved as {fmt} after RGB conversion")
                                        else:
                                            img.save(test_buffer, format=fmt)
                                            logger.info(f"Successfully saved as {fmt}")
                                            
                                    except Exception as fmt_error:
                                        logger.error(f"Failed to save as {fmt}: {str(fmt_error)}")
                                
                                # Count as a successfully processed image
                                image_count += 1
                                logger.info(f"Image {image_count} processed")
                                
                            except Exception as img_error:
                                logger.error(f"Error processing image: {str(img_error)}")
                        
                        except Exception as shape_img_error:
                            logger.error(f"Error processing image shape: {str(shape_img_error)}")
                    
                    # Handle group shapes - contains sub-shapes
                    elif hasattr(shape, 'shape_type') and shape.shape_type == 6:  # GROUP
                        logger.debug("Processing group shape")
                        if hasattr(shape, "shapes"):
                            for i, sub_shape in enumerate(shape.shapes):
                                try:
                                    process_shape(sub_shape, f"{shape_id}_sub{i}")
                                except Exception as sub_error:
                                    logger.error(f"Error processing sub-shape: {str(sub_error)}")
                    
                except Exception as e:
                    logger.error(f"General error processing shape: {str(e)}")
            
            # Process all top-level shapes on the slide
            for i, shape in enumerate(slide.shapes):
                try:
                    process_shape(shape, f"slide{slide_num}_shape{i}")
                except Exception as e:
                    logger.error(f"Top-level shape processing error: {str(e)}")
        
        logger.info(f"Successfully processed {image_count} images")
        return True
        
    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}")
        return False

if __name__ == "__main__":
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        if os.path.exists(filepath):
            process_pptx_images(filepath)
        else:
            logger.error(f"File not found: {filepath}")
    else:
        logger.error("Please provide a PPTX file path")