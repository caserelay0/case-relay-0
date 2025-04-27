import os
import base64
import logging
import io
from typing import Dict, List, Any, Optional
import tempfile
import re
import textwrap
from datetime import datetime
import urllib.parse

# Document processing libraries (already installed)
import PyPDF2
from pdf2image import convert_from_path
import docx
from pptx import Presentation
import trafilatura

from PIL import Image

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

def process_document(filepath_or_url: str) -> Dict[str, Any]:
    """
    Process the uploaded document or web URL and extract text and images
    
    Args:
        filepath_or_url: Path to the uploaded document or a URL
        
    Returns:
        Dictionary containing extracted text and images
    """
    # Set up logging
    logger = logging.getLogger(__name__)
    logger.setLevel(logging.DEBUG)
    
    # Log the file being processed
    logger.debug(f"Processing document: {filepath_or_url}")
    # Check if input is a URL
    if filepath_or_url.startswith(('http://', 'https://')):
        logger.debug(f"Processing as URL: {filepath_or_url}")
        return process_web_content(filepath_or_url)
    
    # Process as a file
    filepath = filepath_or_url
    
    # Get file information
    filename = os.path.basename(filepath)
    file_extension = filename.split('.')[-1].lower() if '.' in filename else ''
    file_stats = os.stat(filepath)
    file_size = file_stats.st_size
    created_time = datetime.fromtimestamp(file_stats.st_ctime)
    
    # Initialize response structure
    result = {
        'text': '',
        'images': [],
        'metadata': {
            'filename': filename,
            'file_type': file_extension,
            'file_size': file_size,
            'created_time': created_time.isoformat(),
            'processing_time': datetime.now().isoformat(),
            'status': 'success',
            'error': None
        }
    }
    
    # Check if file is too large for full processing (general limit for all file types)
    is_large_file = file_size > 15 * 1024 * 1024  # 15MB
    is_very_large_file = file_size > 25 * 1024 * 1024  # 25MB
    
    if is_large_file:
        logger.info(f"Large file detected ({file_size/1024/1024:.1f}MB). Optimizing processing.")
    
    # Flag for automated error recovery
    skip_image_extraction = False
    
    try:
        # Process based on file extension with better error handling
        if file_extension == 'pdf':
            try:
                # For large PDF files, skip image extraction to avoid timeouts
                if is_very_large_file:
                    logger.info(f"Very large PDF file detected ({file_size/1024/1024:.1f}MB). Skipping image extraction.")
                    skip_image_extraction = True
                
                processed = process_pdf(filepath, skip_images=skip_image_extraction)
                if not processed:
                    raise ValueError("Failed to extract content from PDF file")
                
                # Set flags for special handling
                if is_large_file:
                    processed['skip_ai_processing'] = True
                
                result.update(processed)
            except Exception as e:
                logger.error(f"Error processing PDF: {str(e)}")
                result['metadata']['status'] = 'error'
                result['metadata']['error'] = f"PDF processing error: {str(e)}"
                
                # Attempt recovery with text-only extraction
                logger.info("Attempting recovery with text-only extraction")
                try:
                    processed = process_pdf(filepath, skip_images=True)
                    if processed and processed.get('text'):
                        logger.info("Recovery successful with text-only extraction")
                        processed['images'] = []  # Ensure no problematic images
                        processed['skip_ai_processing'] = True  # Skip AI for recovered docs
                        result.update(processed)
                    else:
                        raise ValueError("Recovery failed - could not extract text")
                except Exception as recovery_err:
                    logger.error(f"Recovery failed: {str(recovery_err)}")
                    raise ValueError(f"Failed to process PDF file: {str(e)}")
                
        elif file_extension in ['doc', 'docx']:
            try:
                processed = process_docx(filepath)
                if not processed:
                    raise ValueError("Failed to extract content from Word document")
                
                # Set flags for special handling of large files
                if is_large_file:
                    processed['skip_ai_processing'] = True
                    
                    # For very large files, skip images
                    if is_very_large_file:
                        logger.info(f"Very large DOCX file detected ({file_size/1024/1024:.1f}MB). Skipping images.")
                        processed['images'] = []
                        
                result.update(processed)
            except Exception as e:
                logger.error(f"Error processing DOCX: {str(e)}")
                result['metadata']['status'] = 'error'
                result['metadata']['error'] = f"Word document processing error: {str(e)}"
                raise ValueError(f"Failed to process Word document: {str(e)}")
                
        elif file_extension == 'pptx':
            try:
                # For PPTX files, conditionally skip image extraction
                skip_images_for_pptx = is_very_large_file
                
                # Process the PPTX with options
                processed = process_pptx(filepath, skip_images=skip_images_for_pptx)
                if not processed:
                    raise ValueError("Failed to extract content from PowerPoint presentation")
                
                # Check if file is larger than threshold and set flags for special processing
                if is_large_file:
                    logger.info(f"Large PPTX file detected ({file_size/1024/1024:.1f}MB). Setting optimization flags.")
                    # Skip AI processing for large files
                    processed['skip_ai_processing'] = True
                    
                    # For very large files, also skip image extraction to prevent errors
                    if is_very_large_file:
                        logger.info(f"Very large PPTX file detected ({file_size/1024/1024:.1f}MB). Skipping image extraction.")
                        # Replace existing images with an empty array to avoid processing errors
                        processed['images'] = []
                    
                result.update(processed)
            except Exception as e:
                logger.error(f"Error processing PPTX: {str(e)}")
                result['metadata']['status'] = 'error'
                result['metadata']['error'] = f"PowerPoint processing error: {str(e)}"
                
                # Attempt recovery with text-only extraction
                logger.info("Attempting recovery with text-only extraction for PPTX")
                try:
                    processed = process_pptx(filepath, skip_images=True)
                    if processed and processed.get('text'):
                        logger.info("Recovery successful with text-only extraction for PPTX")
                        processed['images'] = []  # Ensure no problematic images
                        processed['skip_ai_processing'] = True  # Skip AI for recovered docs
                        result.update(processed)
                    else:
                        raise ValueError("Recovery failed - could not extract text from PPTX")
                except Exception as recovery_err:
                    logger.error(f"PPTX recovery failed: {str(recovery_err)}")
                    raise ValueError(f"Failed to process PowerPoint presentation: {str(e)}")
                
        elif file_extension == 'txt':
            # Simple text file handling with multiple encoding attempts
            try:
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    text_content = f.read()
            except UnicodeDecodeError:
                # Try different encodings
                logger.warning("Unicode decode error with utf-8, trying alternative encodings")
                for encoding in ['latin-1', 'iso-8859-1', 'windows-1252']:
                    try:
                        with open(filepath, 'r', encoding=encoding) as f:
                            text_content = f.read()
                            logger.info(f"Successfully read file with encoding: {encoding}")
                            break
                    except Exception:
                        continue
                else:
                    # If we get here, all encodings failed
                    err_msg = "Failed to decode text file with any encoding"
                    logger.error(err_msg)
                    result['metadata']['status'] = 'error'
                    result['metadata']['error'] = err_msg
                    raise ValueError(err_msg)
            structured_content = extract_structured_content(text_content, 'txt')
            result.update({
                'text': text_content,
                'structured_content': structured_content
            })
        else:
            error_msg = f"Unsupported file extension: {file_extension}"
            logger.error(error_msg)
            result['metadata']['status'] = 'error'
            result['metadata']['error'] = error_msg
    except Exception as e:
        error_msg = f"Error processing document: {str(e)}"
        logger.error(error_msg)
        result['metadata']['status'] = 'error'
        result['metadata']['error'] = error_msg
        
        # Provide more specific error information based on error type
        if 'PDF' in str(e) or 'pdf' in str(e):
            raise ValueError(f"PDF processing error: {str(e)}. Your PDF file may be corrupted, password-protected, or in an unsupported format.")
        elif 'Word' in str(e) or 'docx' in str(e) or 'DOCX' in str(e):
            raise ValueError(f"Word document error: {str(e)}. Your document may be corrupted or in an unsupported format.")
        elif 'PowerPoint' in str(e) or 'pptx' in str(e) or 'PPTX' in str(e):
            raise ValueError(f"PowerPoint presentation error: {str(e)}. Your presentation may be corrupted or in an unsupported format.")
        elif 'Image' in str(e) or 'image' in str(e) or 'img' in str(e) or 'IMG' in str(e):
            raise ValueError(f"Image processing error: {str(e)}. Some images in your document could not be processed.")
        elif 'memory' in str(e).lower() or 'Memory' in str(e):
            raise MemoryError(f"Not enough memory to process this document: {str(e)}. Try with a smaller file.")
        elif 'timeout' in str(e).lower() or 'time' in str(e).lower():
            raise TimeoutError(f"Processing timed out: {str(e)}. The document may be too large or complex.")
        elif 'connection' in str(e).lower() or 'reset' in str(e).lower() or 'broken' in str(e).lower():
            raise ConnectionError(f"Connection issue: {str(e)}. The processing was interrupted due to a network or server problem.")
        else:
            raise ValueError(f"Document processing error: {str(e)}. Please try with a different file format or a smaller file.")
    
    # Basic text analysis (word count, etc.)
    if result['text']:
        word_count = len(re.findall(r'\b\w+\b', result['text']))
        result['metadata']['word_count'] = word_count
        result['metadata']['page_count'] = len(result['images']) if file_extension == 'pdf' else None
    
    return result

def process_pdf(filepath: str, skip_images: bool = False) -> Dict[str, Any]:
    """
    Process PDF document
    
    Args:
        filepath: Path to the PDF file
        skip_images: If True, skip image extraction to prevent errors and timeouts
        
    Returns:
        Dictionary containing extracted text and images
    """
    logger.debug(f"Processing PDF: {filepath}")
    
    # Extract text
    text_content = ""
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text_content += page.extract_text() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        text_content = ""
    
    # Extract images - both embedded images and page renders
    images = []
    embedded_images_found = False
    
    # Skip image extraction if requested
    if skip_images:
        logger.info("Skipping image extraction for PDF file as requested")
        return {
            'text': text_content,
            'images': [],
            'structured_content': extract_structured_content(text_content, 'pdf')
        }
    
    # First try to extract embedded images directly from PDF
    try:
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            image_index = 0
            
            # Iterate through each page looking for XObject images
            for page_num, page in enumerate(pdf_reader.pages):
                if '/Resources' in page and '/XObject' in page['/Resources']:
                    xobject = page['/Resources']['/XObject']
                    
                    # Check if XObject is a reference or direct dictionary
                    if hasattr(xobject, "get_object"):
                        xobject = xobject.get_object()
                    
                    # Extract each image from the XObject dictionary
                    for obj in xobject:
                        if xobject[obj]['/Subtype'] == '/Image':
                            try:
                                # Access image data
                                size = (xobject[obj]['/Width'], xobject[obj]['/Height'])
                                
                                # Handle different stream encodings
                                if '/FlateDecode' in xobject[obj]['/Filter']:
                                    # Extract data for supported formats
                                    img_data = xobject[obj]._data
                                    if img_data:
                                        img_buffer = io.BytesIO(img_data)
                                        try:
                                            # Try to open as image
                                            img = Image.open(img_buffer)
                                            img_format = img.format.lower() if img.format else 'jpeg'
                                            img_buffer = io.BytesIO()
                                            img.save(img_buffer, format=img_format.upper())
                                            img_buffer.seek(0)
                                            
                                            # Encode as base64
                                            img_base64 = base64.b64encode(img_buffer.read()).decode('utf-8')
                                            
                                            # Add to image list
                                            images.append({
                                                'id': f"pdf_embedded_{image_index}",
                                                'type': img_format,
                                                'data': img_base64,
                                                'caption': f"Embedded image {image_index+1} (Page {page_num+1})"
                                            })
                                            
                                            image_index += 1
                                            embedded_images_found = True
                                            logger.debug(f"Extracted embedded image {image_index} from page {page_num+1}")
                                        except Exception as img_error:
                                            logger.warning(f"Error processing embedded image: {str(img_error)}")
                            except Exception as obj_error:
                                logger.warning(f"Error extracting embedded object: {str(obj_error)}")
                                continue
                
        logger.debug(f"Embedded image extraction complete. Found {image_index} images.")
    except Exception as e:
        logger.error(f"Error during embedded image extraction: {str(e)}")
    
    # If no embedded images found or extraction failed, fall back to page rendering
    if not embedded_images_found:
        try:
            # Try with a lower dpi if we have issues
            try:
                logger.debug(f"Converting PDF to images: {filepath}")
                # Add a short timeout to prevent hanging on problematic PDFs
                pdf_images = convert_from_path(filepath, dpi=150, timeout=60)
                logger.debug(f"Successfully converted PDF to {len(pdf_images)} images")
            except Exception as e:
                logger.error(f"Error converting PDF to images with default settings: {str(e)}")
                # Try with different settings
                logger.debug("Trying with different settings...")
                try:
                    pdf_images = convert_from_path(filepath, dpi=100, use_pdftocairo=True, timeout=60)
                    logger.debug(f"Successfully converted PDF with alternative settings: {len(pdf_images)} images")
                except Exception as e2:
                    logger.error(f"Error with alternative conversion method: {str(e2)}")
                    # Skip image extraction but continue with text
                    return {
                        'text': text_content,
                        'images': images,  # Return any embedded images we may have found
                        'structured_content': extract_structured_content(text_content, 'pdf')
                    }
            
            # Process each page image
            for i, img in enumerate(pdf_images):
                try:
                    # Save image to buffer
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format='JPEG', quality=85)  # Reduce quality to decrease size
                    img_buffer.seek(0)
                    
                    # Encode image as base64
                    img_base64 = base64.b64encode(img_buffer.read()).decode('utf-8')
                    
                    # Add image to list
                    images.append({
                        'id': f"pdf_page_{i}",
                        'type': 'jpeg',
                        'data': img_base64,
                        'caption': f"Page {i+1}"
                    })
                    logger.debug(f"Processed page {i+1}")
                except Exception as page_error:
                    logger.error(f"Error processing page {i+1}: {str(page_error)}")
                    # Continue with other pages
                    continue
        except Exception as e:
            logger.error(f"Critical error extracting images from PDF: {str(e)}")
            # Continue without images
    
    # Extract structured content
    structured_content = extract_structured_content(text_content, 'pdf')
    
    return {
        'text': text_content,
        'images': images,
        'structured_content': structured_content
    }

def process_docx(filepath: str) -> Dict[str, Any]:
    """
    Process DOCX document
    
    Args:
        filepath: Path to the DOCX file
        
    Returns:
        Dictionary containing extracted text and images
    """
    logger.debug(f"Processing DOCX: {filepath}")
    
    # Extract text
    text_content = ""
    images = []
    
    try:
        doc = docx.Document(filepath)
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            text_content += para.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_content += cell.text + "\n"
        
        # Extract images
        temp_dir = tempfile.mkdtemp()
        image_index = 0
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    img_buffer = io.BytesIO(image_data)
                    
                    # Try to open the image to verify it's valid
                    try:
                        with Image.open(img_buffer) as img:
                            img_format = img.format.lower()
                    except:
                        continue
                    
                    img_buffer.seek(0)
                    img_base64 = base64.b64encode(img_buffer.read()).decode('utf-8')
                    
                    images.append({
                        'id': f"docx_image_{image_index}",
                        'type': img_format,
                        'data': img_base64,
                        'caption': f"Image {image_index+1}"
                    })
                    
                    image_index += 1
                except Exception as e:
                    logger.error(f"Error extracting image from DOCX: {str(e)}")
    except Exception as e:
        logger.error(f"Error processing DOCX: {str(e)}")
    
    # Extract structured content
    structured_content = extract_structured_content(text_content, 'docx')
    
    return {
        'text': text_content,
        'images': images,
        'structured_content': structured_content
    }

def extract_structured_content(text: str, doc_type: str) -> Dict[str, Any]:
    """
    Extract structured content from text based on document type
    
    Args:
        text: The text content of the document
        doc_type: The type of document (pdf, docx, pptx)
        
    Returns:
        Dictionary containing structured information extracted from the text
    """
    # Initialize structured content
    structured = {
        'title': None,
        'summary': None,
        'sections': [],
        'key_points': [],
        'entities': {
            'organizations': [],
            'people': [],
            'dates': []
        }
    }
    
    if not text or len(text.strip()) < 10:
        return structured
    
    # Extract potential title (first non-empty line or first heading)
    lines = text.split('\n')
    non_empty_lines = [line.strip() for line in lines if line.strip()]
    if non_empty_lines:
        structured['title'] = non_empty_lines[0]
        
    # Identify sections using regex patterns for headings
    section_patterns = [
        r'^#+\s+(.+)$',  # Markdown headings
        r'^(\d+\.[\d\.]*\s+.+)$',  # Numbered headings (1.1, 1.2.3, etc.)
        r'^(Chapter \d+:?.*)$',  # Chapter headings
        r'^(.*:)$',  # Colon headings
        r'^([A-Z][A-Z\s]+)$'  # UPPERCASE HEADINGS
    ]
    
    current_section = {'title': 'Introduction', 'content': ''}
    
    for line in lines:
        is_heading = False
        for pattern in section_patterns:
            if re.match(pattern, line.strip()):
                # If we found a new heading, save the current section and start a new one
                if current_section['content'].strip():
                    structured['sections'].append(current_section)
                current_section = {'title': line.strip(), 'content': ''}
                is_heading = True
                break
        
        if not is_heading:
            current_section['content'] += line + '\n'
    
    # Add the last section
    if current_section['content'].strip():
        structured['sections'].append(current_section)
    
    # Extract dates
    date_patterns = [
        r'\d{1,2}/\d{1,2}/\d{2,4}',  # MM/DD/YYYY or DD/MM/YYYY
        r'\d{1,2}-\d{1,2}-\d{2,4}',  # MM-DD-YYYY or DD-MM-YYYY
        r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]* \d{1,2},? \d{4}\b'  # Month DD, YYYY
    ]
    
    for pattern in date_patterns:
        dates = re.findall(pattern, text)
        structured['entities']['dates'].extend(dates)
    
    # Extract potential organizations (capitalized terms)
    org_pattern = r'\b([A-Z][A-Za-z]+ (?:Inc|LLC|Ltd|Corporation|Corp|Company|Co|Group|Partners|Technologies|Solutions|Systems|Associates)\b)'
    organizations = re.findall(org_pattern, text)
    structured['entities']['organizations'] = list(set(organizations))
    
    # Extract potential names (based on common prefixes)
    name_pattern = r'\b(?:Mr|Ms|Mrs|Dr|Prof)\. ([A-Z][a-z]+ [A-Z][a-z]+)\b'
    names = re.findall(name_pattern, text)
    structured['entities']['people'] = list(set(names))
    
    # Generate key points (sections with short content)
    for section in structured['sections']:
        section_content = section['content'].strip()
        if 10 < len(section_content) < 200:  # Short, meaningful sections
            key_point = section_content.replace('\n', ' ')
            structured['key_points'].append(key_point)
    
    # Make sure we have at least a few key points
    if len(structured['key_points']) < 3:
        # Take the first sentence from longer sections
        for section in structured['sections']:
            if len(section['content']) > 200:
                sentences = re.split(r'(?<=[.!?])\s+', section['content'])
                if sentences and len(sentences[0]) > 10:
                    structured['key_points'].append(sentences[0])
                if len(structured['key_points']) >= 5:
                    break
    
    # Limit the number of key points
    structured['key_points'] = structured['key_points'][:7]
    
    return structured

def process_web_content(url: str) -> Dict[str, Any]:
    """
    Process web content and extract text using trafilatura
    
    Args:
        url: URL of the webpage to process
        
    Returns:
        Dictionary containing extracted text and metadata
    """
    logger.debug(f"Processing web content from URL: {url}")
    
    # We've already confirmed trafilatura is installed
    
    # Extract text
    text_content = ""
    images = []
    
    try:
        # Parse URL to get domain and filename-like info
        parsed_url = urllib.parse.urlparse(url)
        domain = parsed_url.netloc
        path = parsed_url.path.strip('/')
        filename = path.split('/')[-1] if path else domain
        
        # Download content
        downloaded = trafilatura.fetch_url(url)
        if downloaded:
            # Extract main content
            text_content = trafilatura.extract(downloaded, include_comments=False, 
                                            include_tables=True, 
                                            output_format='text')
            
            # Extract metadata
            metadata = trafilatura.extract_metadata(downloaded)
            
            # Title and date if available
            title = metadata.title if metadata and metadata.title else None
            date = metadata.date if metadata and metadata.date else None
            
            # Extract structured content
            structured_content = extract_structured_content(text_content, 'web')
            
            # Add title if available
            if title and not structured_content['title']:
                structured_content['title'] = title
            
            return {
                'text': text_content,
                'images': images,  # No images yet, would need to process them separately
                'structured_content': structured_content,
                'metadata': {
                    'url': url,
                    'domain': domain,
                    'title': title,
                    'date': date,
                    'status': 'success'
                }
            }
        else:
            return {
                'text': '',
                'images': [],
                'structured_content': extract_structured_content('', 'web'),
                'metadata': {
                    'url': url,
                    'status': 'error',
                    'error': 'Failed to download content'
                }
            }
    except Exception as e:
        logger.error(f"Error processing web content: {str(e)}")
        return {
            'text': '',
            'images': [],
            'structured_content': extract_structured_content('', 'web'),
            'metadata': {
                'url': url,
                'status': 'error',
                'error': str(e)
            }
        }

def process_pptx(filepath: str, skip_images: bool = False) -> Dict[str, Any]:
    """
    Process PPTX document
    
    Args:
        filepath: Path to the PPTX file
        skip_images: If True, skip image extraction to prevent errors and timeouts
        
    Returns:
        Dictionary containing extracted text and images
    """
    logger.debug(f"Processing PPTX: {filepath}")
    
    text_content = ""
    images = []
    
    # Constants for limiting data size
    MAX_IMAGES = 100  # Maximum number of images to extract
    MAX_SLIDES_PER_CHUNK = 30  # Maximum number of slides to process in one chunk
    
    try:
        presentation = Presentation(filepath)
        total_slides = len(presentation.slides)
        logger.debug(f"PPTX contains {total_slides} slides")
        
        # Skip image extraction if requested
        if skip_images:
            logger.info("Skipping image extraction for PPTX file as requested")
            # Complete text extraction but skip image extraction
            for chunk_start in range(0, total_slides, MAX_SLIDES_PER_CHUNK):
                chunk_end = min(chunk_start + MAX_SLIDES_PER_CHUNK, total_slides)
                logger.debug(f"Processing slide chunk {chunk_start+1} to {chunk_end}")
                
                # Extract text from slides in this chunk
                for i in range(chunk_start, chunk_end):
                    slide = presentation.slides[i]
                    text_content += f"Slide {i+1}:\n"
                    
                    # Get title if present
                    if slide.shapes.title:
                        text_content += f"Title: {slide.shapes.title.text}\n"
                    
                    # Extract text from all shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            # Check if this is not the title we already added
                            if not (slide.shapes.title and shape == slide.shapes.title):
                                text_content += shape.text + "\n"
                    
                    text_content += "\n"
            
            # Return early with text content but no images
            structured_content = extract_structured_content(text_content, 'pptx')
            return {
                'text': text_content,
                'images': [],
                'structured_content': structured_content
            }
        
        # Process in chunks to avoid memory/timeout issues
        for chunk_start in range(0, total_slides, MAX_SLIDES_PER_CHUNK):
            chunk_end = min(chunk_start + MAX_SLIDES_PER_CHUNK, total_slides)
            logger.debug(f"Processing slide chunk {chunk_start+1} to {chunk_end}")
            
            # Extract text from slides in this chunk
            for i in range(chunk_start, chunk_end):
                slide = presentation.slides[i]
                text_content += f"Slide {i+1}:\n"
                
                # Get title if present
                if slide.shapes.title:
                    text_content += f"Title: {slide.shapes.title.text}\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Check if this is not the title we already added
                        if not (slide.shapes.title and shape == slide.shapes.title):
                            text_content += shape.text + "\n"
                
                text_content += "\n"
        
        # Extract images - limit total count and process strategically
        image_index = 0
        
        # Determine which slides to process for images
        # For very large presentations, focus on important slides (first few, last few, and key slides)
        if total_slides > 50:
            # Process first 10 slides, last 10 slides, and every 5th slide in between
            slides_to_process = list(range(0, min(10, total_slides)))
            slides_to_process.extend(range(max(0, total_slides-10), total_slides))
            slides_to_process.extend(range(10, total_slides-10, 5))
            # Remove duplicates and sort
            slides_to_process = sorted(set(slides_to_process))
            logger.debug(f"Large presentation detected. Processing {len(slides_to_process)} slides for images out of {total_slides} total")
        else:
            # Process all slides for smaller presentations
            slides_to_process = range(total_slides)
        
        # Process slide images in two ways:
        # 1. Extract images directly from shapes
        # 2. Also look for images in groups and other containers
        
        for slide_index in slides_to_process:
            # Stop if we've reached the maximum image count
            if image_index >= MAX_IMAGES:
                logger.debug(f"Reached maximum image count ({MAX_IMAGES}). Stopping image extraction.")
                break
                
            slide = presentation.slides[slide_index]
            logger.debug(f"Processing images for slide {slide_index+1}")
            slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_index+1}"
            
            # Helper function to process shapes recursively (including in groups)
            def process_shape(shape, parent_title=None):
                nonlocal image_index
                
                # Stop if we've reached the maximum image count
                if image_index >= MAX_IMAGES:
                    return
                
                # Process based on shape type
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        img_buffer = io.BytesIO(image.blob)
                        
                        # Try to determine image format
                        try:
                            with Image.open(img_buffer) as img:
                                img_format = img.format.lower() if img.format else 'jpeg'
                                img_width, img_height = img.size
                                
                                # Skip very small images (likely icons or bullets)
                                if img_width < 50 or img_height < 50:
                                    logger.debug(f"Skipping small image ({img_width}x{img_height}) on slide {slide_index+1}")
                                    return
                                
                                # Resize large images to reduce data size and handle problematic formats
                                try:
                                    # Convert problematic image formats to JPEG to avoid issues
                                    if img_format.lower() not in ['jpeg', 'jpg', 'png']:
                                        logger.debug(f"Converting image from {img_format} to JPEG format")
                                        if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                            # If image has transparency, convert to PNG
                                            background = Image.new('RGBA', img.size, (255, 255, 255, 255))
                                            background.paste(img, mask=img.split()[3])
                                            img = background.convert('RGB')
                                            img_format = 'jpeg'
                                        else:
                                            # Otherwise convert to JPEG
                                            img = img.convert('RGB')
                                            img_format = 'jpeg'

                                    if img_width > 1000 or img_height > 1000:
                                        logger.debug(f"Resizing large image ({img_width}x{img_height}) on slide {slide_index+1}")
                                        ratio = min(1000 / img_width, 1000 / img_height)
                                        new_width = int(img_width * ratio)
                                        new_height = int(img_height * ratio)
                                        
                                        # Use a safer resampling method
                                        try:
                                            img = img.resize((new_width, new_height), Image.LANCZOS)
                                        except Exception:
                                            # Fallback to a simpler resampling method if LANCZOS fails
                                            img = img.resize((new_width, new_height), Image.BICUBIC)
                                    
                                    # Save the processed image to a new buffer with error handling
                                    img_buffer = io.BytesIO()
                                    
                                    # Use clean save method with robust error handling
                                    try:
                                        if img_format.lower() in ['jpg', 'jpeg']:
                                            img.save(img_buffer, format='JPEG', quality=75, optimize=True)
                                        else:
                                            img.save(img_buffer, format='PNG', optimize=True)
                                    except Exception as save_err:
                                        logger.warning(f"Error saving image, trying fallback format: {str(save_err)}")
                                        # Always fallback to JPEG on error
                                        img = img.convert('RGB')
                                        img.save(img_buffer, format='JPEG', quality=70)
                                        img_format = 'jpeg'
                                        
                                    img_buffer.seek(0)
                                    
                                except Exception as img_process_err:
                                    logger.error(f"Error processing/resizing image: {str(img_process_err)}")
                                    # Continue with original image buffer on error
                                    img_buffer.seek(0)
                        except Exception as img_err:
                            logger.warning(f"Error determining image format: {str(img_err)}")
                            img_format = 'jpeg'
                        
                        # Get image description if available
                        img_alt_text = ""
                        if hasattr(shape, "alt_text") and shape.alt_text:
                            img_alt_text = shape.alt_text
                        
                        img_buffer.seek(0)
                        img_base64 = base64.b64encode(img_buffer.read()).decode('utf-8')
                        
                        # Determine a good caption based on available information
                        caption = img_alt_text if img_alt_text else f"Image from {parent_title or slide_title}"
                        
                        images.append({
                            'id': f"pptx_image_{image_index}",
                            'type': img_format,
                            'data': img_base64,
                            'caption': caption
                        })
                        
                        image_index += 1
                        logger.debug(f"Extracted image {image_index} from slide {slide_index+1}")
                    except Exception as e:
                        logger.error(f"Error extracting image from PPTX shape: {str(e)}")
                
                # Process group shapes (which may contain images)
                elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    try:
                        if hasattr(shape, "shapes"):
                            group_title = shape.text if hasattr(shape, "text") and shape.text else parent_title
                            for sub_shape in shape.shapes:
                                try:
                                    if image_index < MAX_IMAGES:
                                        process_shape(sub_shape, group_title)
                                except Exception as sub_err:
                                    logger.warning(f"Error processing sub-shape: {str(sub_err)}")
                                    continue  # Skip problematic sub-shapes and continue with the next one
                    except Exception as group_err:
                        logger.warning(f"Error processing group shape: {str(group_err)}")
                
                # Look for embedded images in other shape types (limit to first few slides for efficiency)
                elif slide_index < 20 and hasattr(shape, "fill") and hasattr(shape.fill, "type"):
                    # Check if the shape has a picture fill
                    try:
                        if shape.fill.type == 6:  # MSO_FILL.PICTURE
                            if hasattr(shape.fill, "fore_color") and hasattr(shape.fill.fore_color, "by_picture"):
                                img_blob = shape.fill.fore_color._blob
                                if img_blob:
                                    img_buffer = io.BytesIO(img_blob)
                                    try:
                                        with Image.open(img_buffer) as img:
                                            img_format = img.format.lower() if img.format else 'jpeg'
                                            
                                            # Handle image processing and format issues
                                            try:
                                                # Convert problematic image formats to JPEG to avoid issues
                                                if img_format.lower() not in ['jpeg', 'jpg', 'png']:
                                                    logger.debug(f"Converting fill image from {img_format} to JPEG format")
                                                    if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                                        # If image has transparency, handle properly
                                                        background = Image.new('RGBA', img.size, (255, 255, 255, 255))
                                                        background.paste(img, mask=img.split()[3])
                                                        img = background.convert('RGB')
                                                        img_format = 'jpeg'
                                                    else:
                                                        # Otherwise convert to JPEG
                                                        img = img.convert('RGB')
                                                        img_format = 'jpeg'
                                                
                                                # Resize large images to reduce data size
                                                if img.width > 1000 or img.height > 1000:
                                                    ratio = min(1000 / img.width, 1000 / img.height)
                                                    new_width = int(img.width * ratio)
                                                    new_height = int(img.height * ratio)
                                                    
                                                    # Use safer resize method with fallback
                                                    try:
                                                        img = img.resize((new_width, new_height), Image.LANCZOS)
                                                    except Exception:
                                                        img = img.resize((new_width, new_height), Image.BICUBIC)
                                                
                                                # Save with error handling
                                                img_buffer = io.BytesIO()
                                                try:
                                                    if img_format.lower() in ['jpg', 'jpeg']:
                                                        img.save(img_buffer, format='JPEG', quality=75, optimize=True)
                                                    else:
                                                        img.save(img_buffer, format='PNG', optimize=True)
                                                except Exception as save_err:
                                                    logger.warning(f"Error saving fill image, using fallback: {str(save_err)}")
                                                    # Fallback to basic JPEG
                                                    img = img.convert('RGB')
                                                    img.save(img_buffer, format='JPEG', quality=70)
                                                    img_format = 'jpeg'
                                                
                                                img_buffer.seek(0)
                                                img_base64 = base64.b64encode(img_buffer.read()).decode('utf-8')
                                            except Exception as img_process_err:
                                                logger.error(f"Error processing fill image: {str(img_process_err)}")
                                                # Skip this image by returning early
                                                return
                                            
                                            images.append({
                                                'id': f"pptx_fill_image_{image_index}",
                                                'type': img_format,
                                                'data': img_base64,
                                                'caption': f"Background image from {parent_title or slide_title}"
                                            })
                                            
                                            image_index += 1
                                            logger.debug(f"Extracted fill image {image_index} from slide {slide_index+1}")
                                    except Exception as img_err:
                                        logger.warning(f"Error processing picture fill: {str(img_err)}")
                    except Exception as fill_err:
                        logger.debug(f"Error checking fill type: {str(fill_err)}")
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                if image_index < MAX_IMAGES:
                    process_shape(shape)
            
            # If no images found on an important slide, note for possible future implementation
            if (slide_index < 3 or slide_index == total_slides - 1) and slide_index not in [s for s, _ in enumerate(images) if s == slide_index]:
                logger.debug(f"No images found on important slide {slide_index+1}")
        
        logger.debug(f"Total images extracted: {len(images)}")
    
    except Exception as e:
        logger.error(f"Error processing PPTX: {str(e)}")
    
    # Extract structured content
    structured_content = extract_structured_content(text_content, 'pptx')
    
    return {
        'text': text_content,
        'images': images,
        'structured_content': structured_content

def split_text(text, max_tokens=5000):
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""
    for p in paragraphs:
        if len(current_chunk) + len(p) < max_tokens:
            current_chunk += p + "\n\n"
        else:
            chunks.append(current_chunk)
            current_chunk = p + "\n\n"
    if current_chunk:
        chunks.append(current_chunk)
    return chunks

